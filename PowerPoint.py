from pptx import Presentation
from copy import deepcopy
from classes import Slide

def populateSlides(contentDict: dict, filePath: str) -> str:
    prs = Presentation(filePath)

    # since we have the template as the first slide (index = 0), we use index 1 as our starting point.
    idx = 1

    while contentDict:
        # create a blank copy based on our template
        slideCopy(prs)
        buildNewSlide(contentDict.get(idx), prs.slides[idx])
        contentDict.pop(idx)
        idx += 1
    
    outputFileName = "PopulatedSlides.pptx"
    prs.save(outputFileName)

    return outputFileName

def buildNewSlide(slideToBuild: Slide, blankSlide: Slide) -> bool:
        for shape in blankSlide.shapes:
            # check for shapes with text frames
            if not shape.has_text_frame:
                continue
            # Text frames are broken into paragraphs and runs.
            # the key for the shape is the value originally in the shape.
            key = ""
            textBox = shape.text_frame
            for paragraph in textBox.paragraphs:
                for run in paragraph.runs:
                    key += run.text
                    run.text = ""
                    # if we have this key in our slide object dictionary, we replace the text.
                    if key in slideToBuild.content_locations:
                        strBuilder = ""
                        for text in slideToBuild.content_locations[key]:
                            strBuilder += text
                            if len(slideToBuild.content_locations[key]) > 1:
                                strBuilder += "\n"
                        run.text = strBuilder
                        slideToBuild.content_locations.pop(key)
        return True


def slideCopy(prs: Presentation) -> Slide:
    slide_to_copy = prs.slides[0]
    slide_layout = prs.slide_layouts.get_by_name("blank")

    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide_to_copy.shapes:
        el = shape.element
        newEl = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newEl, 'p:extLst')
    return new_slide


# def correctWayToClearAndPopulateText():
#    """
#    This would be the correct way to clear and populate the text box, however you lose the inherent formatting.
#    For the purpose of this tool, we want to preserve the users formatting so we just swap text.
#    """
#     textBox.clear()
#     for text in slideToBuild.content_locations[key]:
#         if paragraph.text == "":
#             paragraph.text = text
#         else: 
#             textBox.add_paragraph().text = text
