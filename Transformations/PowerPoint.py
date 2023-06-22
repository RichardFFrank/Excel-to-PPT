from pptx import Presentation
from copy import deepcopy
from Transformations.Classes import Slide
from datetime import datetime


def populateSlides(contentDict: dict, filePath: str) -> str:
    prs = Presentation(filePath)
    
    offset = len(prs.slides)-1
    
    idx = 1

    while contentDict:

        currSlide = contentDict.get(idx)

        # create a blank copy based on our template
        slideCopy(prs, int(currSlide.slide_type))
        
        buildNewSlide(contentDict.get(idx), prs.slides[offset+idx])
        contentDict.pop(idx)
        idx += 1
    
    outputFileName = f"./Assets/populatedslides-{datetime.now().strftime('%d-%m-%Y-%H-%M')}.pptx"
    prs.save(outputFileName)

    return outputFileName


def buildNewSlide(slideToBuild: Slide, blankSlide: Slide) -> bool:
    for shape in blankSlide.shapes:
        # check for shapes with text frames
        if not shape.has_text_frame:
            continue
        # Text frames are broken into paragraphs and runs.
        # the key for the shape is the value originally in the shape.
        key = ""
        textBox = shape.text_frame
        for paragraph in textBox.paragraphs:
            for run in paragraph.runs:
                key += run.text
                run.text = ""
                # if we have this key in our slide object dictionary, we replace the text.
                if key in slideToBuild.content_locations:
                    strBuilder = ""
                    for text in slideToBuild.content_locations[key]:
                        strBuilder += text
                        if len(slideToBuild.content_locations[key]) > 1:
                            strBuilder += "\n"
                    run.text = strBuilder
                    slideToBuild.content_locations.pop(key)
    return True


def slideCopy(prs: Presentation, templateIdx: int) -> Slide:
    slide_to_copy = prs.slides[templateIdx - 1]
    slide_layout = prs.slide_layouts.get_by_name("blank")

    new_slide = prs.slides.add_slide(slide_layout)

    for shape in slide_to_copy.shapes:
        el = shape.element
        newEl = deepcopy(el)
        new_slide.shapes._spTree.insert_element_before(newEl, 'p:extLst')
    return new_slide

# def correctWayToClearAndPopulateText():
#    """
#    This would be the correct way to clear and populate the text box, however you lose the inherent formatting.
#    For the purpose of this tool, we want to preserve the users formatting so we just swap text.
#    """
#     textBox.clear()
#     for text in slideToBuild.content_locations[key]:
#         if paragraph.text == "":
#             paragraph.text = text
#         else: 
#             textBox.add_paragraph().text = text
